
from PySide2.QtCore import QThread, Qt
import asyncio

import csv
from openpyxl.reader import excel
from module._requirement_func import *
from module._translate_j2k import t_j2k

from speech_recognition import AudioFile, Recognizer, UnknownValueError
from navertts import NaverTTS
from pydub import AudioSegment
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import pdf2docx

from tempfile import gettempdir

from PIL import Image, ImageDraw, ImageFont
import uuid
import random
import shutil


tempdir = gettempdir()

ffmpeg_path = os.path.abspath("utils/ffmpeg.exe")


textList = ['.txt', '.srt']
audioList = ['.wav', '.mp3', ]
videoList = ['.mp4', ]
documentList = ['.docx', '.pdf', ]
presentaionList = ['.pptx', ]
excelList = ['.xlsx', '.csv', ]
imageList = ['.jpg', '.png', ]




class FileTranslationThread(QThread):
    def __init__(self, window, chrd_path):
        QThread.__init__(self)
        self.window = window
        self.chrd_path = chrd_path
        self.setTerminationEnabled = True

    def stop(self):
        if self.isRunning():
            self.window.btnSetting(setNum=1)
            self.window.show_status.setText("파일 번역 중지")
            self.terminate()


    async def rft(self, a):
        s = [asyncio.ensure_future(self.runFileTrans(i)) for i in a]
        await asyncio.gather(*s)




    async def runFileTrans(self, file_name:str):
        fileEXT = os.path.splitext(file_name)[-1].lower()
        abs_filePath = self.window.transed_file_dir.text() + '/' + t_j2k(file_name.split('/')[-1].replace(fileEXT, '') + f'_transed{fileEXT}')




        if fileEXT in textList:
            content = ""

            japaneseText = ListChunk(codecs.open(file_name, 'r', encoding='utf-8').read().split('\n'), 3000)

            for i, jt in enumerate(japaneseText):
                start_time = time()
                j = '\n'.join(jt)

                k = t_j2k(
                    japanese=j,
                    isForceFilter=self.window.f_isActivateFF.isChecked(),
                    isRemoveOri=self.window.f_isActivateRO.isChecked(),
                ).replace('? 1', '1')

                if self.window.f_isPrintBoth.isChecked():
                    jp = j.split('\n')
                    kr = k.split('\n')

                    new_k = ""
                    for jap, kor in zip(jp, kr):
                        new_k += jap+'\n'+kor

                else:
                    new_k = k
                
                content += new_k + '\n'
                self.window.f_log_browser.append(f'번역 Chunk {i+1}회 완료 / 걸린 시간 : {time()-start_time}')


            WriteFile(content, abs_filePath)
            
    



        elif fileEXT in videoList:
            # Current Code -> Vrew, 임시로 만들어둔 것 입니다. 이후에 모델이 나올 경우 다른 방법으로 재구성할 예정
            srtFile = os.path.relpath(self.window.transed_file_dir.text() + '\\' + os.path.basename(file_name).split('.')[0] + '.srt')
            try: os.remove(srtFile)
            except FileNotFoundError: pass

            options = Options()
            options.add_argument('--headless')
            options.add_argument("--window-size=1920,1080")
            options.add_experimental_option("excludeSwitches" , ["enable-automation", "load-extension", "enable-logging"])
            options.add_experimental_option('prefs', {'download.default_directory':os.path.abspath(self.window.transed_file_dir.text())})
            options.add_argument(f'--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/92.0.4515.159 Safari/537.36')
            
            chrome_service = ChromeService(self.chrd_path)
            chrome_service.creationflags = CREATE_NO_WINDOW
            driver = Chrome(options=options, service=chrome_service)

            driver.get('https://vrew.voyagerx.com/ko/try/')
            
            sleep(10)
            a = driver.find_element(By.XPATH, '//input[@name="video-chooser"]')
            a.send_keys(os.path.abspath(file_name))
            driver.find_element(By.ID, 'id_ja').click()
            driver.find_element(By.XPATH, '//button[@class="blue-button"]').click()

            sleep(3)
            while True:
                try:
                    driver.find_element(By.CLASS_NAME, 'progress-stt')
                    sleep(2)
                except exceptions.NoSuchElementException:
                    break

            try:
                sleep(3)
                driver.find_element(By.XPATH, '//*[@id="root"]/div/div[1]/div[1]/div[1]/div[1]/button[2]/div[2]').click()
                sleep(3)
                driver.find_element(By.XPATH, '//div[@data-tooltip="다른 형식으로 내보내기"]').click()
                sleep(3)
                driver.find_element(By.XPATH, '//button[@data-tooltip="자막파일로 내보내기 (.srt)"]').click()
                driver.find_element(By.XPATH, '//button[@class="blue-button"]').click()

            except exceptions.ElementNotInteractableException:
                driver.find_element(By.XPATH, '//*[@id="root"]/div/div[1]/div[1]/div[1]/div[1]/button[2]/div[2]').click()
                sleep(3)
                driver.find_element(By.XPATH, '//div[@data-tooltip="다른 형식으로 내보내기"]').click()
                sleep(3)
                driver.find_element(By.XPATH, '//button[@data-tooltip="자막파일로 내보내기 (.srt)"]').click()
                driver.find_element(By.XPATH, '//button[@class="blue-button"]').click()


            
            srtFile = os.path.relpath(self.window.transed_file_dir.text() + '\\' + os.path.basename(file_name).split('.')[0] + '.srt')

            while True:
                if os.path.isfile(srtFile):
                    driver.close()
                    sleep(3)
                    break
            
            transedSRT = t_j2k(codecs.open(srtFile, 'r', encoding='utf-8').read()).replace('? 1', '1')


            WriteFile(transedSRT, srtFile)
        
        
            os.system(f'{ffmpeg_path} -i {os.path.relpath(file_name)} -filter_complex "subtitles={srtFile}:force_style=\'OutlineColour=&000000000,BorderStyle=3,Outline=3,Shadow=0,MarginV=20\'" -c:V libx264 -preset ultrafast -c:a aac {os.path.relpath(abs_filePath)}'.replace('\\', '/').replace('//', '/'))
            os.system(f'del {srtFile}')
            









        elif fileEXT in audioList:
            rec = Recognizer()
            sound = AudioSegment.from_file(file_name)
            sound.export(file_name, format='wav')


            with AudioFile(file_name) as a_src:
                content = rec.record(a_src)
            try:
                j = rec.recognize_google(
                    audio_data=content,
                    language='ja-JP'
                )
            except UnknownValueError:
                return
            
            k = t_j2k(
                    japanese=j,
                    isForceFilter=self.window.f_isActivateFF.isChecked(),
                    isRemoveOri=self.window.f_isActivateRO.isChecked(),
            )

            sound_gender = 'f' if self.window.f_isActive_F_Sound.isChecked() else 'm'

            tts = NaverTTS(k, lang='ko', gender=sound_gender)
            tts.save(abs_filePath)






        elif fileEXT in documentList:
            if fileEXT == '.pdf':
                converted_name = './utils/tmp/converted_{}'.format(random.randint(100000, 999999))
                pdf2docx.parse(file_name, converted_name)
                abs_filePath = abs_filePath.replace('.pdf', '.docx')
                file_name = converted_name
            

            doc = Document(file_name)

            for p in doc.paragraphs:
                for i in p.runs:
                    k = t_j2k(
                        japanese=i.text,
                        isForceFilter=self.window.f_isActivateFF.isChecked(),
                        isRemoveOri=self.window.f_isActivateRO.isChecked(),
                    )
                    i.text = k
                
            for t in doc.tables:
                for row in t.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            k = t_j2k(
                                japanese=paragraph.text,
                                isForceFilter=self.window.f_isActivateFF.isChecked(),
                                isRemoveOri=self.window.f_isActivateRO.isChecked(),
                            )
                            paragraph.text = k

            doc.save(abs_filePath)
            # os.remove(converted_name)








        elif fileEXT in presentaionList:
            prs = Presentation(file_name)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame

                        for para in text_frame.paragraphs:
                            for p in para.runs:
                                k = t_j2k(
                                    japanese=p.text,
                                    isForceFilter=self.window.f_isActivateFF.isChecked(),
                                    isRemoveOri=self.window.f_isActivateRO.isChecked(),
                                )
                                p.text = k

            prs.save(abs_filePath)






        elif fileEXT in excelList:
            if fileEXT.endswith('.csv'):
                rdr = csv.reader(codecs.open(file_name, 'r', encoding='utf-8'))
                wr = csv.writer(codecs.open(abs_filePath, 'w', encoding='utf-8'))

                for line in rdr:
                    k = t_j2k(
                        '\n'.join(line)
                    ).split('\n')
                    wr.writerow(k)
                
                

            else:
                wb = load_workbook(file_name)

                for s in wb.sheetnames:
                    ws = wb[s]

                    for col in ws.columns:
                        for c in col:
                            val = c.value
                            if val != None:
                                k = t_j2k(
                                        japanese=val,
                                        isForceFilter=self.window.f_isActivateFF.isChecked(),
                                        isRemoveOri=self.window.f_isActivateRO.isChecked(),
                                    )
                                c.value = k
                
                wb.save(abs_filePath)
            








        elif fileEXT in imageList:
            base_url = 'https://touhou.ai/imgtrans'
            r = requests.post(f'{base_url}/manual-translate', files={'file':open(file_name, 'rb').read()})
            j = r.json()
            tr = j['trans_result']
            src_text = [res['s'] for res in tr]
            trans_text = t_j2k('\n'.join(src_text)).split('\n')

            a = []
            for s, t in zip(src_text, trans_text):
                a.append({'s': s, 't': t})
            j['trans_result'] = a


            r2 = requests.post(f'{base_url}/post-translation-result', json=j)


            try:
                j = r2.json()
                status = j['status']
                taskId = j['task_id']
                print(j)
                
                img = f'https://touhou.ai/imgtrans/result/{taskId}/final.jpg'

                

                if status == 'successful':
                    ImageDownload(abs_filePath, img)
                        

                if status == 'failed':
                    print("FAILED!")
                    

            except KeyError as ke:
                print(ke)

                            
                        

        else:
            return











        self.window.f_log_browser.append(f'{abs_filePath} 번역되었습니다.')
        self.window.f_log_browser.verticalScrollBar().setValue(self.window.f_log_browser.verticalScrollBar().maximum())
        self.window.count += 1








    def run(self):
        self.window.btnSetting(setNum=0)
        entireFileList = [self.window.SelectedFileList.item(idx).text() for idx in range(self.window.SelectedFileList.count()) if self.window.SelectedFileList.item(idx).checkState() == Qt.Checked]

        if len(entireFileList) == 0:
            self.window.btnSetting(setNum=1)
            return

        self.window.count = 0
        self.window.progressBar.setMaximum(len(entireFileList))

        start_time = time()
        
        async_loop(self.rft, entireFileList)

        self.window.count = len(entireFileList)

        self.window.sec.setText(f"{round(time()-start_time, 3)}초")
        self.window.show_status.setText("파일 번역 성공")
        self.window.btnSetting(setNum=1)